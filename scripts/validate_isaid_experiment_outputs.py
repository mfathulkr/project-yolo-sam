from __future__ import annotations

import argparse
import hashlib
import sys
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"
if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))

from sam3_bbox_study.config import load_config, resolve_path


PIPELINES = {
    "sam3_text": "sam3_text_output_dir",
    "remotesam_text": "remotesam_text_output_dir",
    "segearth_ov3": "segearth_ov3_output_dir",
    "yolo_sam3": "yolo_sam3_output_dir",
    "gt_box_sam3": "gt_box_sam3_output_dir",
    "yolo_sam2": "yolo_sam2_output_dir",
    "grounded_sam2": "grounded_sam2_output_dir",
    "yolo_ringmo_sam": "yolo_ringmo_sam_output_dir",
    "gt_box_ringmo_sam": "gt_box_ringmo_sam_output_dir",
}
STRATA = [
    "no_overlap__low_mask_area",
    "no_overlap__high_mask_area",
    "overlap__low_mask_area",
    "overlap__high_mask_area",
]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Validate final iSAID vehicle experiment artifacts.")
    parser.add_argument("--config", type=Path, default=ROOT / "configs" / "isaid_vehicle_yolo26x_cpu_eval.yaml")
    parser.add_argument("--presentation", type=Path, default=ROOT.parent / "presentation_isaid_vehicle_sam3_sam2_study" / "isaid_vehicle_sam3_sam2_summary.pptx")
    parser.add_argument("--report-dir", type=Path, default=ROOT / "results" / "isaid_vehicle_final_report")
    return parser.parse_args()


def fail(message: str) -> None:
    raise AssertionError(message)


def validate(config: dict, presentation_path: Path) -> list[str]:
    notes: list[str] = []
    prepared_eval = resolve_path(config["paths"]["prepared_dataset_dir"]) / config["dataset"]["eval_split"]
    metadata = pd.read_csv(prepared_eval / "metadata.csv")
    expected_images = len(metadata)
    if expected_images != 128:
        fail(f"Expected 128 eval images, found {expected_images}")
    notes.append(f"Eval images: {expected_images}")

    stratum_counts = metadata["stratum"].value_counts().sort_index()
    for stratum in STRATA:
        count = int(stratum_counts.get(stratum, 0))
        if count != 32:
            fail(f"Expected 32 images in {stratum}, found {count}")
    notes.append("Strata: 32 images in each of the four overlap/mask-area groups")

    for pipeline, output_key in PIPELINES.items():
        mask_dir = resolve_path(config["paths"][output_key]) / "masks"
        raw_dir = resolve_path(config["paths"][output_key]) / "raw"
        mask_count = len(list(mask_dir.glob("*.png")))
        raw_count = len(list(raw_dir.glob("*.json")))
        if mask_count != expected_images:
            fail(f"{pipeline}: expected {expected_images} masks, found {mask_count}")
        if raw_count != expected_images:
            fail(f"{pipeline}: expected {expected_images} raw JSON files, found {raw_count}")
        notes.append(f"{pipeline}: {mask_count} masks and {raw_count} raw JSON files")

    metrics_dir = resolve_path(config["paths"]["sam3_triplet_metrics_dir"])
    per_image = pd.read_csv(metrics_dir / "per_image_stratified_metrics.csv")
    expected_rows = expected_images * len(PIPELINES)
    if len(per_image) != expected_rows:
        fail(f"Expected {expected_rows} per-image metric rows, found {len(per_image)}")
    for stratum in STRATA:
        for pipeline in PIPELINES:
            rows = per_image[(per_image["stratum"] == stratum) & (per_image["pipeline"] == pipeline)]
            if len(rows) != 32:
                fail(f"Expected 32 metric rows for {stratum}/{pipeline}, found {len(rows)}")
    notes.append(f"Per-image metrics: {len(per_image)} rows, complete for all pipelines and strata")

    curated_dir = ROOT / "results" / "isaid_vehicle_visualizations_sam3_triplet_curated"
    curated_path = curated_dir / "selected_curated_samples.csv"
    using_curated = curated_path.exists()
    selected_path = curated_path if using_curated else resolve_path(config["paths"]["sam3_triplet_visualizations_dir"]) / "selected_qualitative_samples.csv"
    selected = pd.read_csv(selected_path)
    expected_selected_per_stratum = 1 if using_curated else 4
    expected_selected = expected_selected_per_stratum * len(STRATA)
    if len(selected) != expected_selected:
        fail(f"Expected {expected_selected} selected qualitative samples, found {len(selected)}")
    selected_counts = selected["stratum"].value_counts()
    for stratum in STRATA:
        if int(selected_counts.get(stratum, 0)) != expected_selected_per_stratum:
            fail(
                f"Expected {expected_selected_per_stratum} selected samples for {stratum}, "
                f"found {int(selected_counts.get(stratum, 0))}"
            )

    area_threshold = float(metadata["mask_area_ratio"].median())
    merged = selected[["file_name", "stratum"]].merge(
        metadata[["file_name", "area_group", "overlap_group", "mask_area_ratio", "max_pair_bbox_iou", "stratum"]],
        on="file_name",
        suffixes=("_selected", "_metadata"),
    )
    for _, row in merged.iterrows():
        expected_stratum = f"{row['overlap_group']}__{row['area_group']}"
        if row["stratum_selected"] != expected_stratum or row["stratum_metadata"] != expected_stratum:
            fail(f"{row['file_name']}: selected stratum does not match metadata")
        if row["overlap_group"] == "no_overlap" and float(row["max_pair_bbox_iou"]) != 0.0:
            fail(f"{row['file_name']}: no_overlap sample has non-zero bbox IoU")
        if row["overlap_group"] == "overlap" and float(row["max_pair_bbox_iou"]) < 0.05:
            fail(f"{row['file_name']}: overlap sample has bbox IoU below 0.05")
        if row["area_group"] == "low_mask_area" and float(row["mask_area_ratio"]) > area_threshold:
            fail(f"{row['file_name']}: low_mask_area sample is above median area threshold")
        if row["area_group"] == "high_mask_area" and float(row["mask_area_ratio"]) < area_threshold:
            fail(f"{row['file_name']}: high_mask_area sample is below median area threshold")
    notes.append(
        f"Selected qualitative samples: {expected_selected_per_stratum} per stratum, "
        f"all validated against median mask-area threshold {area_threshold:.8f}"
    )

    vis_dir = resolve_path(config["paths"]["sam3_triplet_visualizations_dir"])
    for _, row in selected.iterrows():
        if using_curated:
            path = ROOT / str(row["visualization"])
            expected_size = (2820, 1800)
        else:
            path = vis_dir / f"{row['stratum']}__{Path(row['file_name']).stem}.png"
            expected_size = (4500, 750)
        if not path.exists():
            fail(f"Missing visualization: {path}")
        with Image.open(path) as image:
            if image.size != expected_size:
                fail(f"Unexpected visualization size for {path}: {image.size}")
    notes.append(
        f"Qualitative visualization {'cards' if using_curated else 'strips'}: "
        f"all {expected_selected} present at {expected_size[0]}x{expected_size[1]}"
    )

    if not presentation_path.exists():
        fail(f"Missing PowerPoint: {presentation_path}")
    presentation = Presentation(presentation_path)
    if len(presentation.slides) != 9:
        fail(f"Expected 9 PowerPoint slides, found {len(presentation.slides)}")
    picture_counts = [sum(1 for shape in slide.shapes if shape.shape_type == 13) for slide in presentation.slides]
    expected_picture_counts = [1, 1, 1, 1] if using_curated else [4, 4, 4, 4]
    if picture_counts[4:8] != expected_picture_counts:
        fail(f"Expected qualitative image counts {expected_picture_counts} on slides 5-8, found {picture_counts[4:8]}")
    notes.append(f"PowerPoint: 9 slides; qualitative slides contain {expected_selected_per_stratum} image card(s) each")
    return notes


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def write_artifact_manifest(config: dict, report_dir: Path, presentation_path: Path) -> Path:
    metrics_dir = resolve_path(config["paths"]["sam3_triplet_metrics_dir"])
    vis_dir = resolve_path(config["paths"]["sam3_triplet_visualizations_dir"])
    presentation_dir = presentation_path.parent
    paths = [
        report_dir / "REPORT.md",
        report_dir / "QA_MANIFEST.md",
        presentation_path,
    ]
    paths.extend(sorted(presentation_dir.glob("*.pdf")))
    paths.extend(sorted((presentation_dir / "tables").glob("*.csv")))
    paths.extend(sorted((presentation_dir / "figures").glob("*.png")))
    paths.extend(sorted((presentation_dir / "figures" / "sample_cases").glob("*.png")))
    paths.extend(sorted(metrics_dir.glob("*.csv")))
    paths.extend(sorted(vis_dir.glob("*.csv")))
    paths.extend(sorted(vis_dir.glob("*.png")))
    curated_dir = ROOT / "results" / "isaid_vehicle_visualizations_sam3_triplet_curated"
    paths.extend(sorted(curated_dir.glob("*.csv")))
    paths.extend(sorted((curated_dir / "hero_cards").glob("*.png")))

    seen: set[Path] = set()
    rows: list[dict[str, object]] = []
    for path in paths:
        resolved = path.resolve()
        if resolved in seen or not path.exists():
            continue
        seen.add(resolved)
        rows.append(
            {
                "path": str(path),
                "bytes": path.stat().st_size,
                "sha256": sha256(path),
            }
        )
    manifest_path = report_dir / "ARTIFACT_MANIFEST.csv"
    pd.DataFrame(rows).to_csv(manifest_path, index=False)
    return manifest_path


def main() -> None:
    args = parse_args()
    config = load_config(args.config)
    notes = validate(config, args.presentation)

    args.report_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = args.report_dir / "QA_MANIFEST.md"
    lines = [
        "# iSAID Vehicle Experiment QA Manifest",
        "",
        "All checks passed.",
        "",
    ]
    lines.extend(f"- {note}" for note in notes)
    lines.append("")
    lines.append(f"- PowerPoint: `{args.presentation}`")
    lines.append(f"- Config: `{args.config}`")
    manifest_path.write_text("\n".join(lines), encoding="utf-8")
    artifact_manifest_path = write_artifact_manifest(config, args.report_dir, args.presentation)
    print(f"Wrote QA manifest: {manifest_path}")
    print(f"Wrote artifact manifest: {artifact_manifest_path}")


if __name__ == "__main__":
    main()
