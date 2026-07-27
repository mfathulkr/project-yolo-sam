from __future__ import annotations

import argparse
import shutil
import sys
from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

STUDY_ROOT = Path(__file__).resolve().parents[1]
REPO_ROOT = STUDY_ROOT.parents[1]
ROOT = REPO_ROOT
for source_root in (STUDY_ROOT / "src", REPO_ROOT / "src"):
    if str(source_root) not in sys.path:
        sys.path.insert(0, str(source_root))

from yolo_sam.config import load_config, resolve_path


PIPELINE_ORDER = [
    "sam3_text",
    "remotesam_text",
    "segearth_ov3",
    "yolo_sam3",
    "gt_box_sam3",
    "yolo_sam2",
    "grounded_sam2",
    "yolo_ringmo_sam",
    "gt_box_ringmo_sam",
]
STRATUM_ORDER = [
    "no_overlap__low_mask_area",
    "no_overlap__high_mask_area",
    "overlap__low_mask_area",
    "overlap__high_mask_area",
]
STRATUM_LABELS = {
    "no_overlap__low_mask_area": "No overlap / Low area",
    "no_overlap__high_mask_area": "No overlap / High area",
    "overlap__low_mask_area": "Overlap / Low area",
    "overlap__high_mask_area": "Overlap / High area",
}
STRATUM_FILL_COLORS = {
    "No overlap / Low area": "F4F8FF",
    "No overlap / High area": "EEF7F1",
    "Overlap / Low area": "FFF5E8",
    "Overlap / High area": "F7F0FF",
}
PIPELINE_LABELS = {
    "sam3_text": "SAM3 text-only",
    "remotesam_text": "RemoteSAM text",
    "segearth_ov3": "SegEarth-OV3 + SAM3",
    "yolo_sam3": "YOLO + SAM3",
    "gt_box_sam3": "GT bbox + SAM3",
    "yolo_sam2": "YOLO + SAM2",
    "grounded_sam2": "GroundingDINO + SAM2",
    "yolo_ringmo_sam": "YOLO + RingMo-SAM",
    "gt_box_ringmo_sam": "GT bbox + RingMo-SAM",
}
PIPELINE_COLORS = [
    "#4c78a8",
    "#59a14f",
    "#8cd17d",
    "#f58518",
    "#54a24b",
    "#b279a2",
    "#e45756",
    "#72b7b2",
    "#ff9da6",
]
SHORT_PIPELINE_LABELS = {
    "SAM3 text-only": "SAM3 text",
    "RemoteSAM text": "RemoteSAM",
    "SegEarth-OV3 + SAM3": "SegEarth\n+SAM3",
    "YOLO + SAM3": "YOLO+SAM3",
    "GT bbox + SAM3": "GT+SAM3",
    "YOLO + SAM2": "YOLO+SAM2",
    "GroundingDINO + SAM2": "GroundDINO\n+SAM2",
    "YOLO + RingMo-SAM": "YOLO\n+RingMo",
    "GT bbox + RingMo-SAM": "GT\n+RingMo",
}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Write final iSAID vehicle report and PowerPoint summary.")
    parser.add_argument("--config", type=Path, default=STUDY_ROOT / "configs" / "yolo26x_cpu_eval.yaml")
    parser.add_argument(
        "--report-dir",
        type=Path,
        default=STUDY_ROOT / "results" / "pipelines" / "final_report",
    )
    parser.add_argument(
        "--presentation-dir",
        type=Path,
        default=STUDY_ROOT / "reports",
    )
    return parser.parse_args()


def pct(value: float) -> str:
    return f"{100.0 * value:.1f}%"


def metric(value: float) -> str:
    return f"{value:.4f}"


def hms(seconds: float) -> str:
    seconds_int = int(round(seconds))
    hours, rem = divmod(seconds_int, 3600)
    minutes, secs = divmod(rem, 60)
    return f"{hours}h {minutes}m {secs}s"


def read_yolo_training(results_csv: Path) -> tuple[pd.Series, pd.Series]:
    results = pd.read_csv(results_csv)
    results.columns = [column.strip() for column in results.columns]
    best_idx = results["metrics/mAP50-95(B)"].idxmax()
    return results.loc[best_idx], results.iloc[-1]


def ordered(df: pd.DataFrame, column: str, order: list[str]) -> pd.DataFrame:
    result = df.copy()
    result[column] = pd.Categorical(result[column], categories=order, ordered=True)
    return result.sort_values(column).reset_index(drop=True)


def make_overall_plot(overall: pd.DataFrame, output_path: Path) -> None:
    overall = overall.copy()
    overall["pipeline"] = pd.Categorical(overall["pipeline"], PIPELINE_ORDER, ordered=True)
    overall = overall.sort_values("pipeline")
    labels = [SHORT_PIPELINE_LABELS.get(label, label) for label in overall["pipeline_label"].tolist()]
    x = range(len(labels))
    width = 0.2
    fig, ax = plt.subplots(figsize=(10, 5))
    metrics = [
        ("mean_iou", "IoU", "#1f77b4"),
        ("mean_dice", "Dice", "#2ca02c"),
        ("mean_precision", "Precision", "#ff7f0e"),
        ("mean_recall", "Recall", "#d62728"),
    ]
    for offset, (column, label, color) in enumerate(metrics):
        xs = [value + (offset - 1.5) * width for value in x]
        ax.bar(xs, overall[column], width=width, label=label, color=color)
    ax.set_xticks(list(x))
    ax.set_xticklabels(labels, rotation=0, ha="center")
    ax.tick_params(axis="x", labelsize=8)
    ax.set_ylim(0, 1)
    ax.set_ylabel("Mean score")
    ax.set_title("Overall segmentation metrics")
    ax.grid(axis="y", alpha=0.25)
    ax.legend(ncol=4, loc="upper center", bbox_to_anchor=(0.5, 1.12))
    fig.tight_layout()
    fig.savefig(output_path, dpi=180)
    plt.close(fig)


def make_stratum_plot(by_stratum: pd.DataFrame, output_path: Path, metric_column: str, title: str) -> None:
    work = by_stratum.copy()
    work["stratum"] = pd.Categorical(work["stratum"], STRATUM_ORDER, ordered=True)
    work["pipeline"] = pd.Categorical(work["pipeline"], PIPELINE_ORDER, ordered=True)
    pivot = (
        work.sort_values(["stratum", "pipeline"])
        .pivot(index="stratum", columns="pipeline_label", values=metric_column)
        .reindex(STRATUM_ORDER)
    )
    label_order = [PIPELINE_LABELS[pipeline] for pipeline in PIPELINE_ORDER if PIPELINE_LABELS[pipeline] in pivot.columns]
    pivot = pivot.reindex(columns=label_order)
    pivot.index = [STRATUM_LABELS[stratum] for stratum in pivot.index]
    fig, ax = plt.subplots(figsize=(11.5, 6.2))
    pivot.plot(kind="barh", ax=ax, color=PIPELINE_COLORS[: len(pivot.columns)], width=0.72)
    ax.set_xlim(0, 1)
    ax.set_xlabel(metric_column.replace("mean_", "").title())
    ax.set_ylabel("")
    ax.set_title(title)
    ax.tick_params(axis="both", labelsize=11)
    ax.xaxis.label.set_size(11)
    ax.title.set_size(13)
    ax.grid(axis="x", alpha=0.25)
    ax.invert_yaxis()
    ax.legend(loc="lower center", ncol=3, bbox_to_anchor=(0.5, 1.02), frameon=False, fontsize=10)
    fig.subplots_adjust(left=0.22, right=0.98, top=0.82, bottom=0.10)
    fig.tight_layout()
    fig.savefig(output_path, dpi=180)
    plt.close(fig)


def markdown_table(df: pd.DataFrame) -> str:
    headers = [str(column) for column in df.columns]
    lines = [
        "| " + " | ".join(headers) + " |",
        "| " + " | ".join(["---"] * len(headers)) + " |",
    ]
    for row in df.itertuples(index=False):
        values = ["" if pd.isna(value) else str(value) for value in row]
        lines.append("| " + " | ".join(values) + " |")
    return "\n".join(lines)


def write_markdown_report(
    output_path: Path,
    config: dict,
    best_row: pd.Series,
    last_row: pd.Series,
    overall: pd.DataFrame,
    by_stratum: pd.DataFrame,
    pairwise: pd.DataFrame,
    selected: pd.DataFrame,
    area_threshold: float,
) -> None:
    yolo_table = pd.DataFrame(
        [
            {
                "best_epoch": int(best_row["epoch"]),
                "best_time": hms(float(best_row["time"])),
                "stopped_epoch": int(last_row["epoch"]),
                "total_elapsed": hms(float(last_row["time"])),
                "precision": metric(float(best_row["metrics/precision(B)"])),
                "recall": metric(float(best_row["metrics/recall(B)"])),
                "mAP50": metric(float(best_row["metrics/mAP50(B)"])),
                "mAP50-95": metric(float(best_row["metrics/mAP50-95(B)"])),
            }
        ]
    )
    overall_table = overall[
        [
            "pipeline_label",
            "images",
            "mean_iou",
            "mean_dice",
            "mean_precision",
            "mean_recall",
            "mean_pred_area_ratio",
            "zero_iou_images",
        ]
    ].copy()
    for column in ["mean_iou", "mean_dice", "mean_precision", "mean_recall", "mean_pred_area_ratio"]:
        overall_table[column] = overall_table[column].map(metric)

    stratum_table = by_stratum[
        [
            "stratum",
            "pipeline_label",
            "images",
            "mean_objects",
            "mean_mask_area_ratio",
            "mean_iou",
            "mean_dice",
            "mean_precision",
            "mean_recall",
            "zero_iou_images",
        ]
    ].copy()
    for column in ["mean_objects", "mean_mask_area_ratio", "mean_iou", "mean_dice", "mean_precision", "mean_recall"]:
        stratum_table[column] = stratum_table[column].map(metric)

    if "visualization" not in selected.columns:
        selected = selected.copy()
        selected["visualization"] = selected.apply(
            lambda row: f"studies/isaid_vehicle_study/results/pipelines/visualizations_sam3_triplet/{row['stratum']}__{Path(row['file_name']).stem}.png",
            axis=1,
        )
    selected_columns = ["file_name", "stratum", "num_objects", "mask_area_ratio", "max_pair_bbox_iou"]
    if "reason" in selected.columns:
        selected_columns.append("reason")
    selected_columns.append("visualization")
    selected_table = selected[selected_columns].copy()
    for column in ["mask_area_ratio", "max_pair_bbox_iou"]:
        selected_table[column] = selected_table[column].map(metric)

    lines = [
        "# iSAID Vehicle YOLO/SAM Experiment Report",
        "",
        "## Setup",
        "",
        "- Primary dataset for this run: iSAID overhead urban imagery.",
        "- Target object: merged `vehicle` class from iSAID `Small_Vehicle` and `Large_Vehicle` instance polygons.",
        "- GT masks and GT boxes come from iSAID instance annotations; this run does not use Semantic Drone connected components.",
        "- Evaluation split: 128 positive tiles, balanced as 32 images per bbox-overlap / target-mask-area stratum.",
        "- Pipelines: SAM3 text-only, RemoteSAM text/referring, SegEarth-OV3 SAM3 open-vocabulary fusion, YOLO26x + SAM3 box prompt, GT bbox + SAM3 box prompt, YOLO26x + SAM2 box prompt, GroundingDINO + SAM2 zero-shot text-to-box baseline, and remote-sensing fine-tuned RingMo-SAM constrained by YOLO/GT boxes.",
        "- CPU-only inference was used for the final run with `CUDA_VISIBLE_DEVICES=''`.",
        "",
        "## YOLO26x training checkpoint",
        "",
        markdown_table(yolo_table),
        "",
        f"- Best checkpoint: `{config['yolo']['trained_weights']}`",
        "- Training was intentionally stopped and not resumed for this report.",
        "",
        "## Overall segmentation metrics",
        "",
        markdown_table(overall_table),
        "",
        "## Stratified metrics",
        "",
        markdown_table(stratum_table),
        "",
        "## Pairwise IoU deltas by stratum",
        "",
    ]
    if pairwise.empty:
        lines.append("No pairwise table was produced.")
    else:
        pairwise_table = pairwise.copy()
        for column in ["mean_difference", "ci95_low", "ci95_high"]:
            pairwise_table[column] = pairwise_table[column].map(metric)
        lines.append(markdown_table(pairwise_table))

    lines.extend(
        [
            "",
            "## Selected qualitative examples",
            "",
            "- Selection rule: one hand-curated, highly visible urban example per stratum; overlap/no-overlap and low/high area are still validated against eval metadata.",
            f"- Eval median mask-area threshold used for low/high split: `{area_threshold:.8f}`.",
            "",
            markdown_table(selected_table),
            "",
            "## Short interpretation",
            "",
            "- GT bbox + SAM3 has the highest overall IoU, Dice, and recall; RingMo-SAM is more precise but misses too much area, so better detection boxes still help mask quality most.",
            "- YOLO + SAM3 consistently improves over SAM3 text-only in all four strata, so detection guidance is useful even with an imperfect detector.",
            "- RemoteSAM text is the strongest box-free text/referring baseline in this run; it beats SAM3 text-only and SegEarth-OV3 overall while still trailing YOLO + SAM2 and GT bbox + SAM3.",
            "- SegEarth-OV3 improves recall over plain SAM3 text but still over-segments low-mask-area scenes.",
            "- YOLO + SAM2 is slightly better than YOLO + SAM3 overall on this run, especially in low-mask-area strata, but the two are close in overlap-low-mask-area.",
            "- GroundingDINO + SAM2 is included as a CPU-run zero-shot text-to-box baseline, not as a task-finetuned detector.",
            "- RingMo-SAM is included as a remote-sensing fine-tuned model; its optical instance decoder class map is restricted to YOLO or GT boxes.",
            "- SAM3 text-only over-segments heavily in small-object strata, visible in its high predicted-area ratio and lower precision.",
        ]
    )
    output_path.write_text("\n".join(lines), encoding="utf-8")


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.55))
    text_frame = title_box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(25)
    paragraph.font.bold = True
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.48), Inches(0.78), Inches(12.1), Inches(0.35))
        sub_frame = sub_box.text_frame
        sub_frame.clear()
        sub = sub_frame.paragraphs[0]
        sub.text = subtitle
        sub.font.size = Pt(10.5)


def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float, font_size: int = 16) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(font_size)


def hex_to_rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_df_table(
    slide,
    df: pd.DataFrame,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 8,
    group_column: str | None = None,
) -> None:
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    for col_index, column in enumerate(df.columns):
        cell = table.cell(0, col_index)
        cell.text = str(column)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(font_size)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb("E9ECEF")
    for row_index, row in enumerate(df.itertuples(index=False), start=1):
        fill = None
        if group_column and group_column in df.columns:
            group_value = str(df.iloc[row_index - 1][group_column])
            fill = STRATUM_FILL_COLORS.get(group_value)
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(font_size)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(fill)


def add_picture_fit(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def sample_image_path(row: pd.Series, vis_dir: Path) -> Path:
    visualization = row.get("visualization", "")
    if isinstance(visualization, str) and visualization:
        candidate = ROOT / visualization
        if candidate.exists():
            return candidate
    return vis_dir / f"{row['stratum']}__{Path(row['file_name']).stem}.png"


def build_pptx(
    output_path: Path,
    best_row: pd.Series,
    last_row: pd.Series,
    overall: pd.DataFrame,
    by_stratum: pd.DataFrame,
    pairwise: pd.DataFrame,
    selected: pd.DataFrame,
    area_threshold: float,
    figures_dir: Path,
    vis_dir: Path,
) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_title(slide, "iSAID Vehicle: Text SAM vs Detection-Guided SAM")
    add_bullets(
        slide,
        [
            "Goal: compare text-only, detector-guided, oracle-box, zero-shot, and remote-sensing fine-tuned SAM-style variants.",
            "Target: merged vehicle class from iSAID Small_Vehicle + Large_Vehicle.",
            "Eval design: 128 overhead urban tiles, 32 per overlap / mask-area stratum.",
            f"Low/high area split uses eval median mask-area ratio {area_threshold:.6f}.",
            "Final inference was CPU-only; no GPU was used during this reporting run.",
        ],
        left=0.7,
        top=1.25,
        width=11.8,
        height=2.0,
        font_size=20,
    )
    yolo_df = pd.DataFrame(
        [
            ["Best epoch", int(best_row["epoch"])],
            ["Total stopped epoch", int(last_row["epoch"])],
            ["Elapsed training", hms(float(last_row["time"]))],
            ["mAP50 / mAP50-95", f"{metric(float(best_row['metrics/mAP50(B)']))} / {metric(float(best_row['metrics/mAP50-95(B)']))}"],
            ["Precision / Recall", f"{metric(float(best_row['metrics/precision(B)']))} / {metric(float(best_row['metrics/recall(B)']))}"],
        ],
        columns=["YOLO26x checkpoint", "Value"],
    )
    add_df_table(slide, yolo_df, 0.9, 4.0, 7.0, 1.9, font_size=13)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Overall Results", "Mean scores over 128 eval images")
    add_picture_fit(slide, figures_dir / "overall_metrics.png", 0.55, 1.05, 6.8, 3.65)
    overall_ppt = overall[
        ["pipeline_label", "mean_iou", "mean_dice", "mean_precision", "mean_recall"]
    ].copy()
    overall_ppt["pipeline_label"] = overall_ppt["pipeline_label"].map(lambda label: SHORT_PIPELINE_LABELS.get(label, label))
    for column in ["mean_iou", "mean_dice", "mean_precision", "mean_recall"]:
        overall_ppt[column] = overall_ppt[column].map(metric)
    overall_ppt.columns = ["Pipeline", "IoU", "Dice", "Prec.", "Recall"]
    add_df_table(slide, overall_ppt, 7.55, 1.2, 5.15, 2.25, font_size=10)
    add_bullets(
        slide,
        [
            "Detection-guided SAM improves over text-only in every stratum.",
            "GT boxes + SAM3 is the upper-bound detection case.",
            "YOLO + SAM2 is slightly ahead of YOLO + SAM3 overall in this CPU run.",
            "RemoteSAM text is the strongest box-free remote-sensing text/referring baseline.",
            "SegEarth-OV3 tests SAM3 semantic-head + instance-head fusion without boxes.",
            "GroundingDINO + SAM2 is a zero-shot text-to-box baseline, not a fine-tuned detector.",
            "RingMo-SAM adds a remote-sensing fine-tuned box-constrained baseline.",
        ],
        left=7.65,
        top=3.75,
        width=4.8,
        height=2.1,
        font_size=12,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Stratified View", "2x2 split: bbox overlap yes/no x target mask area low/high")
    add_picture_fit(slide, figures_dir / "mean_iou_by_stratum.png", 0.45, 1.05, 6.2, 4.25)
    add_picture_fit(slide, figures_dir / "mean_dice_by_stratum.png", 6.9, 1.05, 6.0, 4.25)
    add_bullets(
        slide,
        [
            "Low-mask-area strata are the hardest: small vehicles punish over-segmentation and missed detections.",
            "High-mask-area strata give all guided methods more signal and higher IoU/Dice.",
        ],
        left=0.8,
        top=5.55,
        width=11.8,
        height=1.0,
        font_size=14,
    )

    compact = by_stratum[["stratum", "pipeline_label", "mean_iou", "mean_dice"]].copy()
    compact["stratum"] = compact["stratum"].astype(str).map(STRATUM_LABELS)
    for column in ["mean_iou", "mean_dice"]:
        compact[column] = compact[column].map(metric)
    compact.columns = ["Stratum", "Pipeline", "IoU", "Dice"]
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Stratified Metrics Table")
    add_df_table(slide, compact, 0.25, 0.85, 12.8, 6.30, font_size=4.7, group_column="Stratum")

    for stratum in STRATUM_ORDER:
        rows = selected[selected["stratum"] == stratum].copy()
        slide = prs.slides.add_slide(blank)
        add_title(
            slide,
            stratum.replace("__", " / ").replace("_", " "),
            "Curated qualitative card: full context, zoomed GT boxes, and model masks.",
        )
        if len(rows) <= 1:
            for _, row in rows.iterrows():
                add_picture_fit(slide, sample_image_path(row, vis_dir), 1.92, 1.13, 9.35, 5.96)
        else:
            for idx, (_, row) in enumerate(rows.head(2).iterrows()):
                top = 1.05 + idx * 3.05
                add_picture_fit(slide, sample_image_path(row, vis_dir), 1.2, top, 10.9, 2.82)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Takeaways")
    add_bullets(
        slide,
        [
            "YOLO guidance matters: YOLO + SAM3 beats SAM3 text-only in every 2x2 stratum.",
            "RemoteSAM text narrows the gap without detector boxes, but does not beat YOLO-guided SAM2.",
            "SegEarth-OV3 adds recall but remains prone to broad masks in low-area strata.",
            "Detection quality matters: GT bbox + SAM3 remains the best IoU/Dice case overall and by stratum.",
            "SAM2 is competitive with the same YOLO boxes and slightly better overall here.",
            "GroundingDINO + SAM2 tests text-to-box prompting separately from YOLO training.",
            "RingMo-SAM tests whether a remote-sensing fine-tuned decoder helps when the box source is fixed.",
            "Text-only SAM3 over-segments small-object scenes, especially low-mask-area tiles.",
            "The selected examples are balanced across overlap/no-overlap and low/high total vehicle mask area.",
        ],
        left=0.7,
        top=1.2,
        width=11.7,
        height=4.0,
        font_size=16,
    )
    prs.save(output_path)


def main() -> None:
    args = parse_args()
    config = load_config(args.config)
    metrics_dir = resolve_path(config["paths"]["sam3_triplet_metrics_dir"])
    vis_dir = resolve_path(config["paths"]["sam3_triplet_visualizations_dir"])

    report_dir = args.report_dir
    presentation_dir = args.presentation_dir
    figures_dir = presentation_dir / "figures"
    tables_dir = presentation_dir / "tables"
    sample_dir = figures_dir / "sample_cases"
    report_dir.mkdir(parents=True, exist_ok=True)
    figures_dir.mkdir(parents=True, exist_ok=True)
    tables_dir.mkdir(parents=True, exist_ok=True)
    if sample_dir.exists():
        shutil.rmtree(sample_dir)
    sample_dir.mkdir(parents=True, exist_ok=True)

    overall = ordered(pd.read_csv(metrics_dir / "summary_overall_stratified.csv"), "pipeline", PIPELINE_ORDER)
    by_stratum = pd.read_csv(metrics_dir / "summary_by_stratum.csv")
    by_stratum["stratum"] = pd.Categorical(by_stratum["stratum"], STRATUM_ORDER, ordered=True)
    by_stratum["pipeline"] = pd.Categorical(by_stratum["pipeline"], PIPELINE_ORDER, ordered=True)
    by_stratum = by_stratum.sort_values(["stratum", "pipeline"]).reset_index(drop=True)
    pairwise = pd.read_csv(metrics_dir / "pairwise_iou_by_stratum.csv")
    curated_dir = (
        STUDY_ROOT
        / "results"
        / "pipelines"
        / "visualizations_sam3_triplet_curated"
    )
    curated_csv = curated_dir / "selected_curated_samples.csv"
    selected = pd.read_csv(curated_csv if curated_csv.exists() else vis_dir / "selected_qualitative_samples.csv")
    selected["stratum"] = pd.Categorical(selected["stratum"], STRATUM_ORDER, ordered=True)
    selected = selected.sort_values(["stratum", "file_name"]).reset_index(drop=True)
    metadata = pd.read_csv(resolve_path(config["paths"]["prepared_dataset_dir"]) / config["dataset"]["eval_split"] / "metadata.csv")
    area_threshold = float(metadata["mask_area_ratio"].median())

    yolo_results_csv = resolve_path(config["paths"]["yolo_train_run_dir"]) / "train" / "results.csv"
    best_row, last_row = read_yolo_training(yolo_results_csv)

    make_overall_plot(overall, figures_dir / "overall_metrics.png")
    make_stratum_plot(by_stratum, figures_dir / "mean_iou_by_stratum.png", "mean_iou", "Mean IoU by stratum")
    make_stratum_plot(by_stratum, figures_dir / "mean_dice_by_stratum.png", "mean_dice", "Mean Dice by stratum")
    make_stratum_plot(by_stratum, figures_dir / "mean_precision_by_stratum.png", "mean_precision", "Mean precision by stratum")
    make_stratum_plot(by_stratum, figures_dir / "mean_recall_by_stratum.png", "mean_recall", "Mean recall by stratum")

    for file_name in [
        "per_image_stratified_metrics.csv",
        "summary_overall_stratified.csv",
        "summary_by_stratum.csv",
        "pairwise_iou_by_stratum.csv",
    ]:
        shutil.copy2(metrics_dir / file_name, tables_dir / file_name)
    selected_table_name = "selected_curated_samples.csv" if curated_csv.exists() else "selected_qualitative_samples.csv"
    shutil.copy2(curated_csv if curated_csv.exists() else vis_dir / "selected_qualitative_samples.csv", tables_dir / selected_table_name)
    if "visualization" in selected.columns:
        for visualization in selected["visualization"]:
            image_path = ROOT / str(visualization)
            if image_path.exists():
                shutil.copy2(image_path, sample_dir / image_path.name)
    else:
        for image_path in sorted(vis_dir.glob("*.png")):
            shutil.copy2(image_path, sample_dir / image_path.name)

    write_markdown_report(
        report_dir / "REPORT.md",
        config,
        best_row,
        last_row,
        overall,
        by_stratum,
        pairwise,
        selected,
        area_threshold,
    )
    build_pptx(
        presentation_dir / "isaid_vehicle_sam3_sam2_summary.pptx",
        best_row,
        last_row,
        overall,
        by_stratum,
        pairwise,
        selected,
        area_threshold,
        figures_dir,
        vis_dir,
    )
    (presentation_dir / "README.md").write_text(
        "\n".join(
            [
                "# iSAID Vehicle SAM3/SAM2 Study",
                "",
                "- `isaid_vehicle_sam3_sam2_summary.pptx`: short presentation deck.",
                "- `isaid_vehicle_sam3_sam2_summary.pdf`: Linux-viewable PDF export of the deck.",
                "- `figures/`: charts and curated qualitative cards.",
                "- `tables/`: CSV metrics copied from the experiment output.",
                f"- Full report: `{report_dir / 'REPORT.md'}`.",
                f"- QA manifest: `{report_dir / 'QA_MANIFEST.md'}`.",
                f"- Artifact manifest: `{report_dir / 'ARTIFACT_MANIFEST.csv'}`.",
            ]
        ),
        encoding="utf-8",
    )
    print(f"Wrote report: {report_dir / 'REPORT.md'}")
    print(f"Wrote deck: {presentation_dir / 'isaid_vehicle_sam3_sam2_summary.pptx'}")


if __name__ == "__main__":
    main()
